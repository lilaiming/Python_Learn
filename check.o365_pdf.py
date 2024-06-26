# -*- coding=utf-8 -*-
# 本脚由lilaiming编写，用于学习使用！
# Email:essid@qq.com

import os
import pandas as pd
from docx import Document
from pptx import Presentation
import PyPDF2
import concurrent.futures

folder_path = os.path.join(os.path.expanduser('~'), 'Desktop', 'output_folder')
commands_file = '2.cfg.cmd.txt'

# 读取要搜索的文本列表
with open(commands_file) as f:
    search_texts = {text.strip() for text in f.readlines()}

missing_commands = {}  # 缺少的命令字典，用于存储文件名和对应的缺少命令列表
all_files = []

# 生成文件列表的生成器
def file_generator(folder_path):
    for file_name in os.listdir(folder_path):
        file_path = os.path.join(folder_path, file_name)
        if os.path.isfile(file_path):
            yield file_path

# 搜索文本并返回找到的命令
def search_text_in_excel(excel_path, search_texts):
    try:
        df = pd.read_excel(excel_path)
        found_commands = search_texts.intersection(set(df.values.flatten()))
        missing_commands[excel_path] = search_texts - found_commands
        return found_commands
    except Exception as e:
        print(f"读取 {excel_path} 时出现错误: {e}")
        return set()

def search_text_in_word(word_path, search_texts):
    try:
        doc = Document(word_path)
        found_commands = set()
        for paragraph in doc.paragraphs:
            for text in search_texts:
                if text in paragraph.text:
                    found_commands.add(text)
        missing_commands[word_path] = search_texts - found_commands
        return found_commands
    except Exception as e:
        print(f"读取 {word_path} 时出现错误: {e}")
        return set()

def search_text_in_ppt(ppt_path, search_texts):
    try:
        prs = Presentation(ppt_path)
        found_commands = set()
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for text in search_texts:
                            if text in paragraph.text:
                                found_commands.add(text)
        missing_commands[ppt_path] = search_texts - found_commands
        return found_commands
    except Exception as e:
        print(f"读取 {ppt_path} 时出现错误: {e}")
        return set()

def search_text_in_pdf(pdf_path, search_texts):
    try:
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            found_commands = set()
            for page in reader.pages:
                text = page.extract_text()
                for search_text in search_texts:
                    if search_text in text:
                        found_commands.add(search_text)
            missing_commands[pdf_path] = search_texts - found_commands
            return found_commands
    except Exception as e:
        print(f"读取 {pdf_path} 时出现错误: {e}")
        return set()

# 处理单个文件
def process_file(file_path):
    file_ext = os.path.splitext(file_path)[1].lower()
    if file_ext in ['.xlsx', '.xls']:
        return search_text_in_excel(file_path, search_texts)
    elif file_ext in ['.docx', '.doc']:
        return search_text_in_word(file_path, search_texts)
    elif file_ext in ['.pptx', '.ppt']:
        return search_text_in_ppt(file_path, search_texts)
    elif file_ext == '.pdf':
        return search_text_in_pdf(file_path, search_texts)
    else:
        return set()

# 生成文件列表
for file_path in file_generator(folder_path):
    all_files.append(file_path)

# 并发处理文件
with concurrent.futures.ThreadPoolExecutor() as executor:
    results = executor.map(process_file, all_files)

# 输出结果
for file_path, found_commands in zip(all_files, results):
    if found_commands:
        filename = os.path.basename(file_path)
        print(f"文件 {filename} 中命令 {', '.join(found_commands)} 出现了 {len(found_commands)} 次")

print("包含全部命令的文件名:")
for file_path, missing_command_set in missing_commands.items():
    if not missing_command_set:
        print(os.path.basename(file_path))

print("缺少命令的文件名:")
for file_path, missing_command_set in missing_commands.items():
    if missing_command_set:
        filename = os.path.basename(file_path)
        print(f"文件 {filename} 缺少命令:")
        for missing_command in missing_command_set:
            print(missing_command)





