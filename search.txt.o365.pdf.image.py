# -*- coding=utf-8 -*-
# 本脚由lilaiming编写，用于学习使用！
# Email:essid@qq.com

import os
import io
import re
import threading
import traceback
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from PIL import Image
import pytesseract
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import pandas as pd
from openpyxl import load_workbook

# 设置Tesseract OCR路径（根据您的实际路径修改）
tesseract_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
pytesseract.pytesseract.tesseract_cmd = tesseract_path


class FileSearcher:
    def __init__(self, root):
        self.root = root
        self.root.title("增强版文件内容搜索工具")
        self.root.geometry("900x700")
        self.root.configure(bg="#f0f0f0")

        # 创建样式
        self.style = ttk.Style()
        self.style.configure("TFrame", background="#f0f0f0")
        self.style.configure("TButton", padding=6, font=("Arial", 10))
        self.style.configure("TLabel", background="#f0f0f0", font=("Arial", 10))
        self.style.configure("TEntry", font=("Arial", 10))
        self.style.configure("Title.TLabel", font=("Arial", 16, "bold"), background="#3f51b5", foreground="white")

        # 变量
        self.folder_path = tk.StringVar()
        self.search_text = tk.StringVar()
        self.tesseract_path = tk.StringVar(value=tesseract_path)
        self.status_text = tk.StringVar(value="就绪")
        self.progress_value = tk.IntVar()
        self.found_count = 0
        self.file_count = 0
        self.processed_count = 0
        self.unmatched_count = 0  # 新增：未匹配文件计数

        # 创建UI
        self.create_widgets()

        # 文件扩展名映射
        self.extensions = {
            '文本文件': ['.txt', '.cfg', '.py', '.log', '.ini', '.bat', '.sh'],
            'Office文档': ['.docx', '.xlsx', '.pptx'],
            'PDF文档': ['.pdf'],
            '图像文件': ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif']
        }

        # 选中的文件类型
        self.selected_types = {
            '文本文件': tk.BooleanVar(value=True),
            'Office文档': tk.BooleanVar(value=True),
            'PDF文档': tk.BooleanVar(value=True),
            '图像文件': tk.BooleanVar(value=True)
        }

        # 创建文件类型选择区域
        self.create_file_type_selector()

    def create_widgets(self):
        # 标题
        title_frame = ttk.Frame(self.root)
        title_frame.pack(fill="x", padx=10, pady=10)
        ttk.Label(title_frame, text="增强版文件内容搜索工具", style="Title.TLabel").pack(fill="x", ipady=10)

        # 主框架
        main_frame = ttk.Frame(self.root)
        main_frame.pack(fill="both", expand=True, padx=10, pady=5)

        # 文件夹选择
        folder_frame = ttk.Frame(main_frame)
        folder_frame.pack(fill="x", pady=5)
        ttk.Label(folder_frame, text="搜索文件夹:").grid(row=0, column=0, padx=5, sticky="w")
        ttk.Entry(folder_frame, textvariable=self.folder_path, width=50).grid(row=0, column=1, padx=5)
        ttk.Button(folder_frame, text="浏览...", command=self.select_folder).grid(row=0, column=2, padx=5)

        # 搜索内容
        search_frame = ttk.Frame(main_frame)
        search_frame.pack(fill="x", pady=5)
        ttk.Label(search_frame, text="搜索内容:").grid(row=0, column=0, padx=5, sticky="w")
        ttk.Entry(search_frame, textvariable=self.search_text, width=50).grid(row=0, column=1, padx=5)

        # Tesseract路径设置
        tesseract_frame = ttk.Frame(main_frame)
        tesseract_frame.pack(fill="x", pady=5)
        ttk.Label(tesseract_frame, text="Tesseract路径:").grid(row=0, column=0, padx=5, sticky="w")
        ttk.Entry(tesseract_frame, textvariable=self.tesseract_path, width=50).grid(row=0, column=1, padx=5)
        ttk.Button(tesseract_frame, text="设置OCR", command=self.set_tesseract_path).grid(row=0, column=2, padx=5)

        # 文件类型选择框架
        self.type_frame = ttk.LabelFrame(main_frame, text="搜索文件类型")
        self.type_frame.pack(fill="x", pady=10)

        # 按钮区域
        button_frame = ttk.Frame(main_frame)
        button_frame.pack(fill="x", pady=10)

        # 修改开始搜索按钮 - 使用深色文字确保清晰可见
        start_button = ttk.Button(
            button_frame,
            text="开始搜索",
            command=self.start_search
        )
        start_button.pack(side="left", padx=5)

        # 为按钮应用自定义样式
        self.style.configure(
            "Start.TButton",
            background="#4CAF50",
            foreground="black",  # 黑色文字确保清晰可见
            font=("Arial", 10, "bold"),  # 加粗字体
            padding=8  # 增加内边距
        )
        start_button.configure(style="Start.TButton")

        ttk.Button(button_frame, text="清除结果", command=self.clear_results).pack(side="left", padx=5)

        # 进度条
        progress_frame = ttk.Frame(main_frame)
        progress_frame.pack(fill="x", pady=5)
        ttk.Label(progress_frame, text="进度:").pack(side="left", padx=5)
        ttk.Progressbar(progress_frame, variable=self.progress_value, maximum=100).pack(side="left", fill="x",
                                                                                        expand=True, padx=5)
        ttk.Label(progress_frame, textvariable=self.status_text).pack(side="left", padx=10)

        # 结果区域
        result_frame = ttk.LabelFrame(self.root, text="搜索结果")
        result_frame.pack(fill="both", expand=True, padx=10, pady=(0, 10))

        # 结果文本框和滚动条
        self.result_text = tk.Text(result_frame, wrap="word", font=("Consolas", 10))
        self.result_text.pack(side="left", fill="both", expand=True, padx=5, pady=5)

        scrollbar = ttk.Scrollbar(result_frame, command=self.result_text.yview)
        scrollbar.pack(side="right", fill="y")
        self.result_text.config(yscrollcommand=scrollbar.set)

        # 配置样式
        self.style.map(
            "Start.TButton",
            background=[('active', '#45a049'), ('!active', '#4CAF50')],
            foreground=[('active', 'black'), ('!active', 'black')]
        )

    def create_file_type_selector(self):
        # 清除现有内容
        for widget in self.type_frame.winfo_children():
            widget.destroy()

        # 创建文件类型选择器
        row, col = 0, 0
        for file_type, var in self.selected_types.items():
            cb = ttk.Checkbutton(
                self.type_frame,
                text=file_type,
                variable=var,
                onvalue=True,
                offvalue=False
            )
            cb.grid(row=row, column=col, padx=10, pady=5, sticky="w")
            col += 1
            if col > 3:
                col = 0
                row += 1

    def select_folder(self):
        folder = filedialog.askdirectory(title="选择要搜索的文件夹")
        if folder:
            self.folder_path.set(folder)

    def set_tesseract_path(self):
        path = self.tesseract_path.get()
        if path and os.path.isfile(path) and path.lower().endswith("tesseract.exe"):
            pytesseract.pytesseract.tesseract_cmd = path
            messagebox.showinfo("成功", "Tesseract OCR路径已设置")
        else:
            messagebox.showerror("错误", "无效的Tesseract路径")

    def start_search(self):
        folder = self.folder_path.get()
        search_text = self.search_text.get().strip()

        if not folder or not os.path.isdir(folder):
            messagebox.showerror("错误", "请选择有效的文件夹路径")
            return

        if not search_text:
            messagebox.showerror("错误", "请输入搜索内容")
            return

        # 重置状态
        self.found_count = 0
        self.file_count = 0
        self.processed_count = 0
        self.unmatched_count = 0  # 重置未匹配计数
        self.progress_value.set(0)
        self.status_text.set("正在统计文件...")
        self.result_text.delete(1.0, tk.END)
        self.result_text.insert(tk.END, "开始搜索...\n")
        self.result_text.see(tk.END)
        self.root.update()

        # 收集要搜索的扩展名
        search_extensions = []
        for file_type, var in self.selected_types.items():
            if var.get():
                search_extensions.extend(self.extensions[file_type])

        if not search_extensions:
            messagebox.showerror("错误", "请至少选择一种文件类型")
            return

        # 在后台线程中执行搜索
        threading.Thread(
            target=self.search_files,
            args=(folder, search_text, search_extensions),
            daemon=True
        ).start()

    def search_files(self, folder, search_text, extensions):
        try:
            # 第一步：统计文件数量
            for root, dirs, files in os.walk(folder):
                for file in files:
                    ext = os.path.splitext(file)[1].lower()
                    if ext in extensions:
                        self.file_count += 1

            if self.file_count == 0:
                self.status_text.set("未找到匹配的文件")
                self.result_text.insert(tk.END, "未找到指定类型的文件\n")
                return

            # 第二步：实际搜索
            self.status_text.set("正在搜索...")
            self.result_text.insert(tk.END, f"总共找到 {self.file_count} 个文件需要搜索\n\n")
            self.root.update()

            # 存储未匹配的文件
            unmatched_files = []

            for root, dirs, files in os.walk(folder):
                for file in files:
                    file_path = os.path.join(root, file)
                    ext = os.path.splitext(file)[1].lower()

                    if ext in extensions:
                        found, content_lines = self.process_file(file_path, ext, search_text)

                        if found:
                            self.found_count += 1
                            self.result_text.insert(tk.END, f"\n在文件中找到匹配: {file_path}\n")
                            for line in content_lines:
                                self.result_text.insert(tk.END, f"  {line}\n")
                            self.result_text.see(tk.END)
                            self.root.update()
                        else:
                            # 记录未匹配的文件
                            unmatched_files.append(file_path)
                            self.unmatched_count += 1

                        self.processed_count += 1
                        progress = int((self.processed_count / self.file_count) * 100)
                        self.progress_value.set(progress)
                        self.status_text.set(f"处理中: {self.processed_count}/{self.file_count} 文件")
                        self.root.update()

            # 完成搜索 - 显示匹配结果统计
            self.result_text.insert(tk.END, f"\n\n{'=' * 80}\n")
            self.result_text.insert(tk.END, f"搜索完成! 在 {self.found_count} 个文件中找到匹配内容\n")

            # 显示未匹配文件列表
            if unmatched_files:
                self.result_text.insert(tk.END, f"\n{'=' * 80}\n")
                self.result_text.insert(tk.END, f"未匹配的文件 (共 {self.unmatched_count} 个):\n\n")

                # 按文件名排序
                unmatched_files.sort()

                # 分组显示以提高可读性
                files_per_line = 3
                for i in range(0, len(unmatched_files), files_per_line):
                    line_files = unmatched_files[i:i + files_per_line]
                    line_text = "    ".join([os.path.basename(f) for f in line_files])
                    self.result_text.insert(tk.END, f"  {line_text}\n")

                self.result_text.insert(tk.END, f"\n未匹配文件路径列表:\n")
                for file_path in unmatched_files:
                    # 显示相对路径以缩短显示
                    rel_path = os.path.relpath(file_path, folder)
                    self.result_text.insert(tk.END, f"  {rel_path}\n")
            else:
                self.result_text.insert(tk.END, f"\n所有文件都包含匹配内容!\n")

            # 显示最终统计
            self.result_text.insert(tk.END, f"\n{'=' * 80}\n")
            self.result_text.insert(tk.END, f"最终统计:\n")
            self.result_text.insert(tk.END, f"  匹配文件: {self.found_count} 个\n")
            self.result_text.insert(tk.END, f"  未匹配文件: {self.unmatched_count} 个\n")
            self.result_text.insert(tk.END, f"  总计文件: {self.file_count} 个\n")

            # 更新状态栏
            self.status_text.set(f"搜索完成! 匹配: {self.found_count} 文件, 未匹配: {self.unmatched_count} 文件")
            self.result_text.see(tk.END)

        except Exception as e:
            self.result_text.insert(tk.END, f"\n发生错误: {str(e)}\n")
            self.result_text.insert(tk.END, traceback.format_exc())
            self.result_text.see(tk.END)
            self.status_text.set("搜索出错")

    def process_file(self, file_path, ext, search_text):
        try:
            found = False
            content_lines = []

            # 处理文本文件
            if ext in self.extensions['文本文件']:
                found, content_lines = self.search_text_file(file_path, search_text)

            # 处理Word文档
            elif ext == '.docx':
                found, content_lines = self.search_docx(file_path, search_text)

            # 处理Excel文档
            elif ext == '.xlsx':
                found, content_lines = self.search_xlsx(file_path, search_text)

            # 处理PPT文档
            elif ext == '.pptx':
                found, content_lines = self.search_pptx(file_path, search_text)

            # 处理PDF文档
            elif ext == '.pdf':
                found, content_lines = self.search_pdf(file_path, search_text)

            # 处理图像文件
            elif ext in self.extensions['图像文件']:
                found, content_lines = self.search_image(file_path, search_text)

            return found, content_lines

        except Exception as e:
            # 处理出错时视为未匹配
            self.result_text.insert(tk.END, f"\n处理文件 {file_path} 时出错: {str(e)}\n")
            return False, []

    def search_text_file(self, file_path, search_text):
        found = False
        content_lines = []

        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                lines = f.readlines()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='gbk') as f:
                    lines = f.readlines()
            except Exception:
                return False, []
        except Exception:
            return False, []

        for line_num, line in enumerate(lines, 1):
            if search_text in line:
                found = True
                # 高亮匹配文本
                highlighted = line.replace(search_text, f"**{search_text}**")
                content_lines.append(f"行 {line_num}: {highlighted.strip()}")

        return found, content_lines

    def search_docx(self, file_path, search_text):
        found = False
        content_lines = []

        try:
            doc = docx.Document(file_path)
            for para_num, para in enumerate(doc.paragraphs, 1):
                if search_text in para.text:
                    found = True
                    # 高亮匹配文本
                    highlighted = para.text.replace(search_text, f"**{search_text}**")
                    content_lines.append(f"段落 {para_num}: {highlighted.strip()}")
        except Exception:
            return False, []

        return found, content_lines

    def search_xlsx(self, file_path, search_text):
        found = False
        content_lines = []

        try:
            wb = load_workbook(filename=file_path, read_only=True)
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell and search_text in str(cell):
                            found = True
                            # 高亮匹配文本
                            cell_value = str(cell).replace(search_text, f"**{search_text}**")
                            content_lines.append(f"工作表 '{sheet_name}': 单元格值: {cell_value}")
                            break
        except Exception:
            return False, []

        return found, content_lines

    def search_pptx(self, file_path, search_text):
        found = False
        content_lines = []

        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if search_text in shape.text:
                            found = True
                            # 高亮匹配文本
                            highlighted = shape.text.replace(search_text, f"**{search_text}**")
                            content_lines.append(f"幻灯片 {slide_num}: {highlighted.strip()}")
        except Exception:
            return False, []

        return found, content_lines

    def search_pdf(self, file_path, search_text):
        found = False
        content_lines = []

        try:
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()

                if search_text in text:
                    found = True

                    # 提取包含搜索词的上下文
                    lines = text.split('\n')
                    for line_num, line in enumerate(lines, 1):
                        if search_text in line:
                            # 高亮匹配文本
                            highlighted = line.replace(search_text, f"**{search_text}**")
                            content_lines.append(f"页面 {page_num + 1}, 行 {line_num}: {highlighted.strip()}")
        except Exception:
            return False, []

        return found, content_lines

    def search_image(self, file_path, search_text):
        found = False
        content_lines = []

        try:
            # 使用Pillow打开图像
            img = Image.open(file_path)

            # 使用Tesseract OCR识别文本
            text = pytesseract.image_to_string(img)

            if search_text in text:
                found = True

                # 提取包含搜索词的上下文
                lines = text.split('\n')
                for line_num, line in enumerate(lines, 1):
                    if search_text in line:
                        # 高亮匹配文本
                        highlighted = line.replace(search_text, f"**{search_text}**")
                        content_lines.append(f"行 {line_num}: {highlighted.strip()}")
        except Exception as e:
            return False, []

        return found, content_lines

    def clear_results(self):
        self.result_text.delete(1.0, tk.END)
        self.progress_value.set(0)
        self.status_text.set("就绪")
        self.found_count = 0
        self.file_count = 0
        self.processed_count = 0
        self.unmatched_count = 0  # 重置未匹配计数


if __name__ == "__main__":
    root = tk.Tk()
    app = FileSearcher(root)
    root.mainloop()

